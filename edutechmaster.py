from flask import Flask, render_template, request,send_file,jsonify,flash,url_for
from gtts import gTTS
import openai
import PIL.Image
import textwrap
import os
import google.generativeai as genai
from datetime import datetime
from bs4 import BeautifulSoup
from docx import Document
import requests
import json
from pptx import Presentation
from flask_wtf import FlaskForm
from wtforms import StringField, PasswordField, SubmitField
from wtforms.validators import DataRequired, Email, EqualTo
from flask_sqlalchemy import SQLAlchemy
from flask_bcrypt import Bcrypt
from sqlalchemy.exc import IntegrityError
from IPython.display import display, Markdown
app = Flask(__name__)
app.config['SECRET_KEY'] = 'your_secret_key'

app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///site.db'
db = SQLAlchemy(app)
bcrypt = Bcrypt(app)



class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(20), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password = db.Column(db.String(60), nullable=False)

class RegistrationForm(FlaskForm):
    username = StringField('Username', validators=[DataRequired()])
    email = StringField('Email', validators=[DataRequired(), Email()])
    password = PasswordField('Password', validators=[DataRequired()])
    confirm_password = PasswordField('Confirm Password', validators=[DataRequired(), EqualTo('password')])
    submit = SubmitField('Sign Up')

class LoginForm(FlaskForm):
    email = StringField('Email', validators=[DataRequired(), Email()])
    password = PasswordField('Password', validators=[DataRequired()])
    submit = SubmitField('Login')


@app.route("/register", methods=['GET', 'POST'])
def register():
    form = RegistrationForm()
    if form.validate_on_submit():
        hashed_password = bcrypt.generate_password_hash(form.password.data).decode('utf-8')
        user = User(username=form.username.data, email=form.email.data, password=hashed_password)
        try:
            db.session.add(user)
            db.session.commit()
            flash('Your account has been created! You are now able to log in', 'success')
            return redirect(url_for('login'))
        except IntegrityError:
            db.session.rollback()
            flash('An account with this email address already exists. Please use a different email address.', 'danger')
    return render_template('register.html', title='Register', form=form)

@app.route("/login", methods=['GET', 'POST'])
def login():
    form = LoginForm()
    if form.validate_on_submit():
        user = User.query.filter_by(email=form.email.data).first()
        if user and bcrypt.check_password_hash(user.password, form.password.data):
            flash('You have been logged in!', 'success')
            return redirect(url_for('index'))
        else:
            flash('Login Unsuccessful. Please check email and password', 'danger')
    return render_template('login.html', title='Login', form=form)



openai.api_key = "Your-api-key"

# Configure Generative AI API key
genai.configure(api_key="Your-api-key")

# Initialize Generative AI model
model = genai.GenerativeModel('gemini-pro-vision')

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/speech')
def speechindex():
    return render_template('speech.html')

@app.route('/store_transcription', methods=['POST'])
def store_transcription():
    data = request.json
    transcription = data['transcription']

    # Store transcription in a Word file
    document = Document()
    document.add_paragraph(transcription)
    document.save('transcription.docx')

    return 'Transcription stored successfully'



@app.route('/speech-summary')
def get_summary():
    # Read the text from the .docx file
    docx_file = 'speechtext.docx'
    doc_text = read_docx(docx_file)

    # Generate summary
    summary = generate_summary(doc_text)

    return summary


def read_docx(file_path):
    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)


def to_markdown(text):
    text = text.replace('\u2022', '  *')
    return Markdown(textwrap.indent(text, '> ', predicate=lambda _: True))


def generate_summary(note_text):
    model = genai.GenerativeModel('gemini-pro')
    rply = model.generate_content("Here is my speech text for an presentation please look for grammer mistakes and where should i improve and also give me the improved speech text with perfect grammer " + note_text)
    return render_template('transcript-result.html', script_Summary=rply.text)



@app.route('/ai-writer')
def writer():
    return render_template('letter.html')
	
@app.route('/generate-letter', methods=['POST'])
def generate_letter():
    prompt = request.form['prompt']
    model = genai.GenerativeModel('gemini-pro')
    rply = model.generate_content(" Generate an " + prompt + " with proper alignment ")
    letter = rply.text

    document = Document()
    document.add_paragraph(letter)
    document.save('letter.docx')
    
    return render_template('index.html')

@app.route('/Edu')
def edu():
    # Format date in day-month-year format
    for item in news_data:
        # Parse the ISO 8601 formatted date and convert it to day-month-year
        item['date'] = datetime.fromisoformat(item['date'].replace('Z', '+00:00')).strftime('%d-%m-%Y')

    return render_template('news.html', news=news_data)

# API endpoint for news data
url = 'https://oevortex-webscout.hf.space/api/news'
params = {
    'q': 'Technology',  # Query parameter
    'max_results': 10,   # Maximum number of results
    'safesearch': 'moderate',  # Safe search option
    'region': 'wt-wt'   # Region parameter
}

headers = {
    'accept': 'application/json'
}

# Get news data from API
response = requests.get(url, params=params, headers=headers)

if response.status_code == 200:
    news_data = response.json()['results']  # Access 'results' key
    # or do something with the data
else:
    print("Error:", response.status_code)


def generate_response_with_format(prompt):
    # Generate the completion
    response = openai.Completion.create(
        model="gpt-3.5-turbo-instruct",  # Use the model you prefer
        prompt=prompt,
        max_tokens=1500,  
        stop=None,  
        temperature=0.5,  
        n=1,  
        echo=False  
    )
    return response.choices[0].text.strip()

@app.route('/ppt', methods=['GET', 'POST'])
def ppt():
    if request.method == 'POST':
        title = request.form['title']
        prompt = f"Generate content slides for {title} presentation. Each slide should have a title and content. The content should cover various features and aspects of the platform. Here is the format for your response: [('title1', 'content1'), ('title2', 'content2'), ...]"
        
        # Generate response from GPT-3
        response = generate_response_with_format(prompt)
        
        
            # Parse the response as JSON
        slides_data = json.loads(response)
        

        # Create a presentation
        prs = Presentation()

        # Add content slides
        for slide_title, content in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            shapes = slide.shapes
            title_shape = shapes.title
            content_shape = shapes.placeholders[1]

            title_shape.text = slide_title
            content_shape.text = content

        # Save the presentation
        filename = f"{title}_Presentation.pptx"
        prs.save(filename)

        return send_file(filename, as_attachment=True)

    return render_template('ppt.html')



def generate_question(note_text):
    model = genai.GenerativeModel('gemini-pro')
    rply = model.generate_content("Generate important qeustions from the given notes and study material,list me the questions"+note_text)
    to_markdown(rply.text)
    return rply.text
# Route for home page with file upload form


@app.route('/important', methods=['GET', 'POST'])
def important():
    if request.method == 'POST':
        # Check if a file was uploaded
        if 'Note_file' not in request.files:
            return render_template('error.html', message='No file part')

        file = request.files['Note_file']

        # Check if the file is empty
        if file.filename == '':
            return render_template('error.html', message='No selected file')

        # Check if the file is of allowed type
        if file and file.filename.endswith('.txt'):
            # Read the file content
            note_text = file.read().decode('utf-8')

            impquestions = generate_question(note_text)
            
       
            return render_template('questionn-result.html', impquestions=impquestions)

        else:
            return render_template('error.html', message='Invalid file type. Please upload a text file')

    return render_template('question.html')



def generate_notesummary(note_text):
    model = genai.GenerativeModel('gemini-pro')
    rply = model.generate_content("summarize my notes"+note_text)
    to_markdown(rply.text)
    return rply.text
# Route for home page with file upload form
@app.route('/note', methods=['GET', 'POST'])
def note():
    if request.method == 'POST':
        # Check if a file was uploaded
        if 'Note_file' not in request.files:
            return render_template('error.html', message='No file part')

        file = request.files['Note_file']

        # Check if the file is empty
        if file.filename == '':
            return render_template('error.html', message='No selected file')

        # Check if the file is of allowed type
        if file and file.filename.endswith('.txt'):
            # Read the file content
            note_text = file.read().decode('utf-8')

            # Generate summary
            summary_text = generate_notesummary(note_text)
            
            # Render the result template with summary
            return render_template('note-result.html', summary_text=summary_text)

        else:
            return render_template('error.html', message='Invalid file type. Please upload a text file')

    return render_template('note.html')



@app.route('/web')
def web():
    return render_template('web-summary.html')


def summarize_text(text):
    model = genai.GenerativeModel('gemini-pro')
    rply = model.generate_content("Summarize this webpage and give me the summary of this webpage"+text)
    to_markdown(rply.text)
    return rply.text


@app.route('/web_summary', methods=['POST'])
def web_summary():
        url = request.form['url']
        websummary = summarize_text(url)
        tts = gTTS(text=websummary, lang='en')
        output_file_audio = "static/outputsummary.mp3"
        tts.save(output_file_audio)
        return render_template('web-result.html', websummary=websummary, audio_file=output_file_audio)
   


@app.route('/chatbot')
def chatbot():
    return render_template('chatbot.html')

@app.route('/get_response', methods=['POST'])
def get_response():
    data = request.json
    user_message = data['message']
    response = openai.Completion.create(
        engine="gpt-3.5-turbo-instruct",
        prompt=user_message,
        max_tokens=150
    )
    gpt_response = response.choices[0].text.strip()
    return jsonify({'response': gpt_response})

    


if __name__ == '__main__':
    app.run(debug=True)